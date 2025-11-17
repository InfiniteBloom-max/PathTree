import pdfplumber
from pptx import Presentation
import re
from typing import Dict, List, Any
import os

class DocumentProcessor:
    def __init__(self):
        self.chunk_size = 1000
        self.overlap = 200
    
    async def process_document(self, file_path: str) -> Dict[str, Any]:
        """Process uploaded document and extract text"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return await self._process_pdf(file_path)
        elif file_extension == '.pptx':
            return await self._process_pptx(file_path)
        elif file_extension == '.txt':
            return await self._process_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    async def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract text from PDF"""
        text_content = ""
        page_count = 0
        
        try:
            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content += page_text + "\n\n"
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")
        
        chunks = self._chunk_text(text_content)
        
        return {
            "raw_text": text_content,
            "chunks": chunks,
            "page_count": page_count,
            "file_type": "pdf"
        }
    
    async def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract text from PowerPoint"""
        text_content = ""
        slide_count = 0
        
        try:
            prs = Presentation(file_path)
            slide_count = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                text_content += slide_text + "\n\n"
                
        except Exception as e:
            raise Exception(f"Error processing PPTX: {str(e)}")
        
        chunks = self._chunk_text(text_content)
        
        return {
            "raw_text": text_content,
            "chunks": chunks,
            "page_count": slide_count,
            "file_type": "pptx"
        }
    
    async def _process_txt(self, file_path: str) -> Dict[str, Any]:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text_content = file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                text_content = file.read()
        except Exception as e:
            raise Exception(f"Error processing TXT: {str(e)}")
        
        chunks = self._chunk_text(text_content)
        
        return {
            "raw_text": text_content,
            "chunks": chunks,
            "page_count": 1,
            "file_type": "txt"
        }
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into chunks for processing"""
        # Clean text
        text = re.sub(r'\s+', ' ', text).strip()
        
        if len(text) <= self.chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + self.chunk_size
            
            if end >= len(text):
                chunks.append(text[start:])
                break
            
            # Try to break at sentence boundary
            sentence_end = text.rfind('.', start, end)
            if sentence_end > start:
                end = sentence_end + 1
            else:
                # Try to break at word boundary
                word_end = text.rfind(' ', start, end)
                if word_end > start:
                    end = word_end
            
            chunks.append(text[start:end])
            start = end - self.overlap
        
        return chunks
    
    def extract_key_phrases(self, text: str) -> List[str]:
        """Extract key phrases from text"""
        # Simple keyword extraction
        words = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', text)
        
        # Remove duplicates and filter
        key_phrases = list(set(words))
        key_phrases = [phrase for phrase in key_phrases if len(phrase) > 3]
        
        return key_phrases[:20]  # Return top 20